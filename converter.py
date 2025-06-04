import streamlit as st
import io
import os
from docx import Document
from fpdf import FPDF
import pandas as pd
from pptx import Presentation

def converter_tool():
    st.title("🔄 File Converter")
    file = st.file_uploader("Upload file to convert", type=["txt", "docx", "pdf", "xlsx", "pptx"])

    if file:
        file_ext = file.name.split(".")[-1].lower()

        # TXT to PDF
        if file_ext == "txt":
            text = file.read().decode("utf-8")
            pdf = FPDF()
            pdf.add_page()
            pdf.set_auto_page_break(auto=True, margin=15)
            pdf.set_font("Arial", size=12)
            for line in text.split("\n"):
                pdf.multi_cell(0, 10, line)
            output = io.BytesIO()
            pdf.output(output)
            st.success("Converted TXT to PDF.")
            st.download_button("⬇️ Download PDF", data=output.getvalue(), file_name="converted.pdf")

        # DOCX to TXT
        elif file_ext == "docx":
            doc = Document(file)
            full_text = "\n".join([p.text for p in doc.paragraphs])
            st.success("Converted DOCX to TXT.")
            st.download_button("⬇️ Download TXT", data=full_text, file_name="converted.txt")

        # XLSX to CSV
        elif file_ext == "xlsx":
            df = pd.read_excel(file)
            csv_data = df.to_csv(index=False)
            st.success("Converted XLSX to CSV.")
            st.download_button("⬇️ Download CSV", data=csv_data, file_name="converted.csv")

        # PPTX to TXT (slide text only)
        elif file_ext == "pptx":
            prs = Presentation(file)
            all_text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        all_text += shape.text + "\n"
            st.success("Converted PPTX to TXT.")
            st.download_button("⬇️ Download TXT", data=all_text, file_name="slides.txt")

        else:
            st.warning("Conversion for this file type is not available yet.")
